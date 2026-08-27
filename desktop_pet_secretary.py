"""小綿助 Windows 桌面秘書。

平時只有透明桌寵常駐；左鍵會說打氣話，右鍵選單可開啟秘書首頁。所有本機健康紀錄、
快速記事與尚未同步的任務均保存在使用者 AppData，不寫入網頁 localStorage。
"""

from __future__ import annotations

import json
import os
import random
import re
import shutil
import subprocess
import sys
import threading
import uuid
from datetime import date, datetime, timedelta
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from urllib.parse import urlparse
import tkinter as tk
from tkinter import filedialog, messagebox, ttk

from PIL import Image, ImageGrab

from desktop_pet_preview import ERROR_LOG, DesktopPetPreview


APP_NAME = "小綿助教師秘書"
DATA_ROOT = Path(os.getenv("APPDATA", Path.home())) / "XiaoMianZhuSecretary"
DATA_FILE = DATA_ROOT / "secretary_data.json"
ATTACHMENT_ROOT = DATA_ROOT / "attachments"
DESKTOP_SETTINGS_NAME = "xiaomianzhu_settings.json"
AUTOSTART_VALUE_NAME = "XiaoMianZhuSecretary"
BRIDGE_HOST = "127.0.0.1"
BRIDGE_PORT = 8767
BRIDGE_MAX_BODY = 5 * 1024 * 1024


def bridge_origin_allowed(origin: str, allowed_web_origin: str = "") -> bool:
    """只接受本機預覽與本專案的公開網頁，避免其他網站讀取教師資料。"""
    if not origin or origin == "null":
        return True
    try:
        parsed = urlparse(origin)
        host = (parsed.hostname or "").lower()
    except ValueError:
        return False
    if host in {"127.0.0.1", "localhost"}:
        return True
    if parsed.scheme != "https":
        return False
    configured = str(allowed_web_origin or "").strip().rstrip("/")
    return host == "cona0815.github.io" or bool(configured and origin.rstrip("/") == configured)


class LocalBridgeServer(ThreadingHTTPServer):
    daemon_threads = True
    allow_reuse_address = True


def make_bridge_handler(secretary: "SecretaryPet") -> type[BaseHTTPRequestHandler]:
    class BridgeHandler(BaseHTTPRequestHandler):
        server_version = "XiaoMianZhuLocalBridge/1.0"

        def log_message(self, _format: str, *_args: object) -> None:
            return

        def _origin(self) -> str:
            return str(self.headers.get("Origin") or "")

        def _headers(self, status: int = 200, content_type: str = "application/json; charset=utf-8") -> bool:
            origin = self._origin()
            if not bridge_origin_allowed(origin, secretary._bridge_allowed_origin()):
                self.send_response(403)
                self.send_header("Content-Type", "application/json; charset=utf-8")
                self.end_headers()
                return False
            self.send_response(status)
            self.send_header("Content-Type", content_type)
            self.send_header("Cache-Control", "no-store")
            self.send_header("Access-Control-Allow-Origin", origin or "null")
            self.send_header("Access-Control-Allow-Methods", "GET, POST, OPTIONS")
            self.send_header("Access-Control-Allow-Headers", "Content-Type")
            self.send_header("Access-Control-Allow-Private-Network", "true")
            self.send_header("Vary", "Origin")
            self.end_headers()
            return True

        def _write_json(self, payload: dict, status: int = 200) -> None:
            if self._headers(status):
                self.wfile.write(json.dumps(payload, ensure_ascii=False).encode("utf-8"))

        def do_OPTIONS(self) -> None:  # noqa: N802 - BaseHTTPRequestHandler API
            self._headers(204)

        def do_GET(self) -> None:  # noqa: N802 - BaseHTTPRequestHandler API
            if self.path == "/health":
                self._write_json(secretary._bridge_health())
                return
            if self.path == "/panel":
                # 秘書面板改由瀏覽器呈現：部分機器的 Tk 文字繪製不穩定（面板空白），
                # 瀏覽器渲染在任何顯示環境都可靠。頁面只在本機回環位址提供。
                if self._headers(200, "text/html; charset=utf-8"):
                    self.wfile.write(secretary._panel_html().encode("utf-8"))
                return
            if self.path == "/panel-data":
                self._write_json(secretary._panel_payload())
                return
            self._write_json({"ok": False, "error": "not found"}, 404)

        def do_POST(self) -> None:  # noqa: N802 - BaseHTTPRequestHandler API
            if self.path not in ("/sync", "/panel-action"):
                self._write_json({"ok": False, "error": "not found"}, 404)
                return
            try:
                length = int(self.headers.get("Content-Length") or "0")
                if length <= 0 or length > BRIDGE_MAX_BODY:
                    raise ValueError("資料大小不正確")
                payload = json.loads(self.rfile.read(length).decode("utf-8"))
                if not isinstance(payload, dict):
                    raise ValueError("資料格式不正確")
                if self.path == "/panel-action":
                    self._write_json(secretary._panel_apply_action(payload))
                    return
                result = secretary._bridge_sync(payload)
                self._write_json(result)
            except (ValueError, json.JSONDecodeError, UnicodeDecodeError) as error:
                self._write_json({"ok": False, "error": str(error)}, 400)

    return BridgeHandler


def now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")


def application_directory() -> Path:
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent
    return Path(__file__).resolve().parent


def external_settings_path() -> Path:
    return application_directory() / DESKTOP_SETTINGS_NAME


def load_external_settings() -> dict:
    path = external_settings_path()
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, ValueError, json.JSONDecodeError):
        return {}
    if payload.get("kind") != "xiaomianzhu-desktop-settings" or not isinstance(payload.get("settings"), dict):
        return {}
    return payload["settings"]


def startup_command() -> str:
    if getattr(sys, "frozen", False):
        return f'"{Path(sys.executable).resolve()}" --startup-launch'
    pythonw = Path(sys.executable).with_name("pythonw.exe")
    runner = pythonw if pythonw.exists() else Path(sys.executable)
    return f'"{runner.resolve()}" "{Path(__file__).resolve()}" --startup-launch'


def configure_windows_autostart(enabled: bool) -> None:
    if os.name != "nt":
        return
    try:
        import winreg

        key_path = r"Software\Microsoft\Windows\CurrentVersion\Run"
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_SET_VALUE) as key:
            if enabled:
                winreg.SetValueEx(key, AUTOSTART_VALUE_NAME, 0, winreg.REG_SZ, startup_command())
            else:
                try:
                    winreg.DeleteValue(key, AUTOSTART_VALUE_NAME)
                except FileNotFoundError:
                    pass
    except OSError:
        # 公司或學校的 Windows 原則可能禁止修改登入啟動項；程式仍可手動執行。
        pass


PANEL_PAGE_HTML = """<!DOCTYPE html>
<html lang="zh-Hant"><head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>__PET_NAME__ 教師秘書</title>
<style>
:root{--bg:#f3f6f4;--surface:#fff;--soft:#e8f2ee;--primary:#1f514a;--strong:#173f3a;--ink:#1c302d;--muted:#657570;--line:#d5e1dc;--danger:#a34339;--danger-soft:#fff0ed;--gold:#b9852f}
*{box-sizing:border-box;margin:0;padding:0}
body{font-family:"Microsoft JhengHei","Noto Sans TC",sans-serif;background:var(--bg);color:var(--ink);font-size:15px;line-height:1.7}
header{background:var(--strong);color:#fff;padding:14px 22px;display:flex;align-items:center;gap:12px}
header h1{font-size:20px}header .meta{margin-left:auto;font-size:12px;opacity:.85}
main{max-width:1180px;margin:0 auto;padding:16px;display:grid;grid-template-columns:1fr 1fr;gap:14px}
@media(max-width:900px){main{grid-template-columns:1fr}}
.card{background:var(--surface);border:1.5px solid var(--line);border-radius:14px;padding:14px 16px}
.card.soft{background:var(--soft)}.card.danger{background:var(--danger-soft)}
.card h2{font-size:16px;color:var(--strong);margin-bottom:8px}
.card ul{list-style:none}.card li{padding:3px 0;border-bottom:1px dashed var(--line)}.card li:last-child{border-bottom:0}
.empty{color:var(--muted);font-size:13px}
button{font-family:inherit;border:0;border-radius:10px;padding:8px 18px;font-size:14px;font-weight:800;cursor:pointer;background:var(--primary);color:#fff}
button.light{background:var(--surface);color:var(--strong);border:1.5px solid var(--line)}
.row{display:flex;gap:8px;align-items:center;flex-wrap:wrap;margin:6px 0}
textarea,input{font-family:inherit;font-size:14px;border:1.5px solid var(--line);border-radius:10px;padding:8px 10px;width:100%}
.tag{display:inline-block;background:var(--soft);color:var(--strong);border-radius:8px;padding:1px 8px;font-size:11px;font-weight:800;margin-right:4px}
.grid2{display:grid;grid-template-columns:1fr 1fr;gap:10px}
.num{font-size:26px;font-weight:900;color:var(--strong)}
.note-meta{color:var(--muted);font-size:11px}
details summary{cursor:pointer;font-weight:800;color:var(--strong)}
.seat-ok{color:#1c6b45}.seat-bad{color:var(--danger);font-weight:800}
</style></head><body>
<header><h1>🐑 __PET_NAME__ 教師秘書</h1><div class="meta">本機頁面（127.0.0.1）｜<span id="stamp">載入中</span>｜每 15 秒自動更新</div></header>
<main>
  <section class="card soft" style="grid-column:1/-1"><h2>☀ 今日簡報</h2><div id="brief" class="empty">載入中……</div></section>

  <section class="card"><h2>📌 今日重要</h2><ul id="dueToday"></ul></section>
  <section class="card danger"><h2>⚠ 已逾期</h2><ul id="overdue"></ul></section>

  <section class="card" style="grid-column:1/-1"><h2>🏫 今日出席與作業繳交（晨間大屏回報）</h2>
    <div id="cloudBody" class="empty">載入中……</div>
    <details id="cloudSetup" style="margin-top:8px"><summary>⚙ 雲端連線設定（第一次使用點我）</summary>
      <div class="row"><input id="cfgUrl" placeholder="Apps Script 部署網址（/exec 結尾，與大屏同一份）"></div>
      <div class="row"><input id="cfgToken" type="password" placeholder="教室大屏金鑰（CLASSROOM_TOKEN）"><input id="cfgClass" placeholder="班級名稱（例：501）" style="max-width:160px"></div>
      <div class="row"><button onclick="saveCloudCfg()">儲存並讀取</button><span class="note-meta">設定只存在這台電腦的瀏覽器。</span></div>
    </details></section>

  <section class="card soft"><h2>🌿 健康管理</h2>
    <div class="grid2">
      <div>💧 今日喝水 <span class="num" id="water">0</span> 杯<div class="row"><button onclick="act('water')">喝一杯</button><span class="note-meta" id="waterHint"></span></div></div>
      <div>🚶 距上次起身 <span class="num" id="moveElapsed">0</span> 分<div class="row"><button onclick="act('move_done')">活動完成</button><span class="note-meta" id="moveHint"></span></div></div>
    </div>
    <div style="margin-top:8px">💊 服藥提醒：<span id="medicine" class="empty">未設定</span></div>
    <details style="margin-top:8px"><summary>⚙ 詳細提醒設定（喝水／起身間隔、服藥時間）</summary>
      <div class="row">💧 喝水提醒間隔
        <select id="cfgWaterInterval" onchange="actx({action:'set_water_interval',minutes:this.value})"><option>30</option><option>45</option><option>60</option><option>90</option><option>120</option></select> 分鐘
        　🚶 起身提醒間隔
        <select id="cfgMoveInterval" onchange="actx({action:'set_move_interval',minutes:this.value})"><option>15</option><option>30</option><option>45</option><option>60</option><option>90</option><option>120</option></select> 分鐘
      </div>
      <div class="row">💊 新增服藥時間 <input id="cfgMedTime" type="time" style="max-width:140px"> <button onclick="addMed()">＋新增</button><span class="note-meta">每一筆時間到了小綿助都會提醒；點時間旁的 ✖ 可移除。</span></div>
      <div class="row" id="medManage"></div>
      <div class="note-meta">提醒的總開關與安靜時段，在教師工作台「設定 → 小綿助」。</div>
    </details>
  </section>

  <section class="card"><h2>📮 待追蹤</h2><ul id="tracking"></ul></section>

  <section class="card"><h2>💬 快速記事</h2>
    <textarea id="noteInput" rows="3" placeholder="先記下臨時交代、等待回覆或稍後要整理的事情……"></textarea>
    <div class="row"><button onclick="addNote()">儲存記事</button><span class="note-meta">保存在小綿助本機資料；工作台開啟時會自動帶入。</span></div>
  </section>

  <section class="card"><h2>📱 LINE 待整理</h2><ul id="lineInbox"></ul><div class="note-meta">請到教師工作台的「LINE 收件匣」確認建立。</div></section>

  <section class="card" style="grid-column:1/-1"><h2>📝 本機記事（最近 10 則）</h2><ul id="notes"></ul></section>

  <section class="card" style="grid-column:1/-1"><h2>🐑 桌寵控制</h2>
    <div class="row">
      <button class="light" onclick="act('toggle_pause')">⏸ 暫停／繼續走動</button>
      <button class="light" style="color:var(--danger);border-color:var(--danger)" onclick="if(confirm('確定關閉小綿助？（提醒與本機連線都會停止）'))act('quit')">❌ 關閉小綿助</button>
      <span class="note-meta">開啟本頁：對桌面上的小綿羊按右鍵或雙擊。</span>
    </div>
  </section>
</main>
<script>
'use strict';
const $=id=>document.getElementById(id);
const esc=s=>String(s).replace(/[&<>"]/g,c=>({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;'}[c]));
function fill(id, items, render, empty){const el=$(id);el.innerHTML=items.length?items.map(render).join(''):'<li class="empty">'+empty+'</li>';}
async function load(){
  try{
    const d=await (await fetch('/panel-data',{cache:'no-store'})).json();
    $('stamp').textContent=d.generatedAt;
    $('brief').innerHTML=`今天有 <b>${d.dueToday.length}</b> 件期限任務、<b>${d.overdue.length}</b> 件逾期、<b>${d.tracking.length}</b> 件等待回覆。`+(d.overdue.length?`　建議優先處理：<b>${esc(d.overdue[0].title)}</b>`:d.dueToday.length?`　今天先完成：<b>${esc(d.dueToday[0].title)}</b>`:'　今天的期限工作已整理完成 🎉');
    fill('dueToday',d.dueToday,t=>`<li>${esc(t.title)}${t.due_time?'｜'+t.due_time:''}</li>`,'今天沒有期限任務');
    fill('overdue',d.overdue,t=>`<li>${esc(t.title)}</li>`,'目前沒有逾期任務');
    fill('tracking',d.tracking,t=>`<li>${esc(t.title)}</li>`,'目前沒有等待回覆的事項');
    fill('lineInbox',d.lineInbox,i=>`<li><span class="tag">${i.kind}</span>${i.medium==='voice'?'🎤':i.medium==='photo'?'📷':''}${esc(i.title)}${i.tag?`<span class="tag">${esc(i.tag)}</span>`:''}</li>`,'目前沒有 LINE 待整理的訊息');
    fill('notes',d.notes,n=>`<li>${esc(n.text)}<div class="note-meta">${esc(n.time)}</div></li>`,'目前沒有記事');
    $('water').textContent=d.health.water;
    $('waterHint').textContent=`（每 ${d.health.waterInterval} 分鐘提醒）`;
    $('moveElapsed').textContent=d.health.moveElapsed;
    $('moveHint').textContent=d.health.moveDone?'（今天已完成 ✅）':`（每 ${d.health.moveInterval} 分鐘提醒）`;
    if(document.activeElement!==$('cfgWaterInterval'))$('cfgWaterInterval').value=String(d.health.waterInterval);
    if(document.activeElement!==$('cfgMoveInterval'))$('cfgMoveInterval').value=String(d.health.moveInterval);
    const med=d.health.medicineTimes;
    $('medicine').innerHTML=med.length?med.map(t=>d.health.medicineDone.includes(t)?`<span class="tag">✅ ${t}</span>`:`<span class="tag">⏰ ${t}</span> <button class="light" style="padding:2px 10px" onclick="act('medicine_done','${t}')">已服用</button>`).join(' '):'未設定';
    $('medManage').innerHTML=med.length?med.map(t=>`<span class="tag">${t} <button class="light" style="padding:0 8px" title="移除" onclick="actx({action:'remove_medicine_time',time:'${t}'})">✖</button></span>`).join(' '):'<span class="note-meta">尚未設定服藥時間。</span>';
  }catch(e){$('brief').textContent='讀取失敗：'+e.message;}
}
async function act(action,time){await fetch('/panel-action',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({action,time})});load();}
async function actx(payload){const r=await(await fetch('/panel-action',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(payload)})).json().catch(()=>null);if(r&&r.ok===false)alert(r.error||'設定失敗');load();}
async function addMed(){const t=$('cfgMedTime').value;if(!t)return alert('請先選擇時間');await actx({action:'add_medicine_time',time:t});$('cfgMedTime').value='';}
async function addNote(){const t=$('noteInput').value.trim();if(!t)return;await fetch('/panel-action',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({action:'add_note',text:t})});$('noteInput').value='';load();}
function cloudCfg(){try{return JSON.parse(localStorage.getItem('petPanelCloud')||'null')||{}}catch(e){return{}}}
function saveCloudCfg(){localStorage.setItem('petPanelCloud',JSON.stringify({url:$('cfgUrl').value.trim(),token:$('cfgToken').value.trim(),cls:$('cfgClass').value.trim()}));loadCloud();}
async function loadCloud(){
  const c=cloudCfg();
  $('cfgUrl').value=c.url||'';$('cfgToken').value=c.token||'';$('cfgClass').value=c.cls||'';
  if(!c.url||!c.token||!c.cls){$('cloudBody').textContent='尚未設定雲端連線——點下方「⚙ 雲端連線設定」貼上大屏用的網址、教室金鑰與班級即可。';$('cloudSetup').open=true;return;}
  $('cloudBody').textContent='讀取大屏回報中……';
  try{
    const r=await(await fetch(c.url,{method:'POST',redirect:'follow',headers:{'Content-Type':'text/plain;charset=utf-8'},body:JSON.stringify({action:'classroom_today',token:c.token,className:c.cls})})).json();
    if(!r.ok)throw new Error(r.error||'雲端回應異常');
    const absent=(r.attendance||[]).map(a=>`${a.seat}號${a.status}`);
    const boards={};
    (r.homework||[]).forEach(h=>{const k=h.subject+(h.assignment?'／'+h.assignment:'');(boards[k]=boards[k]||{m:[],r:[]})[h.status==='補交'?'r':'m'].push(h.seat);});
    let html=`<div>👥 出缺：${absent.length?'<span class="seat-bad">'+absent.join('、')+'</span>':'<span class="seat-ok">全班到齊 🎉</span>'}</div>`;
    const keys=Object.keys(boards);
    html+= keys.length?keys.map(k=>`<div>📚 ${esc(k)}：${boards[k].m.length?'<span class="seat-bad">缺交 '+boards[k].m.join('、')+'號</span>':'<span class="seat-ok">全交 ✅</span>'}${boards[k].r.length?'｜補交 '+boards[k].r.join('、')+'號':''}</div>`).join(''):'<div class="empty">今天還沒有作業登記。</div>';
    $('cloudBody').innerHTML=html;
  }catch(e){$('cloudBody').textContent='大屏回報讀取失敗：'+e.message+'（今天可能還沒送出回報）';}
}
load();loadCloud();setInterval(load,15000);setInterval(loadCloud,60000);
</script></body></html>"""


def sanitize_line_inbox(items: object) -> list[dict]:
    """整理網頁橋接傳來的 LINE 待整理清單；桌面端唯讀顯示，欄位全部截長補短。"""
    if not isinstance(items, list):
        return []
    sanitized: list[dict] = []
    for item in items[:50]:
        if not isinstance(item, dict):
            continue
        title = str(item.get("title") or "").strip()[:80]
        if not title:
            continue
        sanitized.append({
            "id": str(item.get("id") or "")[:60],
            "title": title,
            "type": "task" if item.get("type") == "task" else "note",
            "tag": str(item.get("tag") or "")[:20],
            "medium": item.get("medium") if item.get("medium") in ("voice", "photo") else "",
            "created_at": str(item.get("createdAt") or "")[:40],
        })
    return sanitized


def default_data() -> dict:
    now = now_iso()
    return {
        "tasks": [],
        "notes": [],
        # LINE 待整理清單：由教師工作台網頁經本機橋接同步進來，桌面端唯讀顯示。
        "line_inbox": [],
        "settings": {
            "enabled": True,
            "autostart": False,
            "pet_name": "小綿助",
            "cheer_interval_minutes": 5,
            "startup_delay_seconds": 30,
            "health_reminders": True,
            "opening_brief": True,
            "remember_position": True,
            "quiet_start": "18:00",
            "quiet_end": "07:30",
            "bridge_enabled": True,
            "bridge_last_sync": "",
            "allowed_web_origin": "",
            "device_id": str(uuid.uuid4()),
        },
        "health": {
            "date": date.today().isoformat(),
            "water_count": 0,
            "last_water": now,
            "last_move": now,
            "move_done_date": "",
            "move_interval": 60,
            "water_interval": 60,
            "medicine_time": "",
            "medicine_times": [],
            "medicine_done_times": [],
            "medicine_alerted_times": [],
            "medicine_done_date": "",
        },
    }


class SecretaryPet(DesktopPetPreview):
    CHEER_MESSAGES = [
        "今天也辛苦了，先完成最重要的一件事就很好。",
        "慢慢來，事情一件一件做就會完成。",
        "你的用心，學生和同事都感受得到。",
        "先喝口水，再繼續也不遲。",
        "別忘了替今天努力的自己按個讚。",
        "忙碌時先深呼吸，我陪你整理下一步。",
        "完成一小步，也是一個很棒的進度。",
        "你已經做得很好了，剩下的我們慢慢處理。",
        "今天也要留一點時間照顧自己。",
        "遇到難題沒關係，先把問題拆小一點。",
        "老師加油！你的穩定就是孩子的安心。",
        "休息一下不是偷懶，是為了走得更久。",
    ]

    COLORS = {
        "bg": "#f3f6f4",
        "surface": "#ffffff",
        "soft": "#e8f2ee",
        "primary": "#1f514a",
        "primary_hover": "#173f3a",
        "strong": "#173f3a",
        "ink": "#1c302d",
        "muted": "#657570",
        "line": "#d5e1dc",
        "danger": "#a34339",
        "danger_soft": "#fff0ed",
        "gold": "#b9852f",
        "gold_soft": "#fff7e6",
    }

    def __init__(self) -> None:
        super().__init__()
        self.data = self._load_data()
        self._save_data()
        settings = self.data["settings"]
        configure_windows_autostart(bool(settings.get("enabled", True) and settings.get("autostart", False)))
        if settings.get("remember_position"):
            try:
                self.x = max(0, min(self.screen_width - 80, int(settings.get("pet_x", self.x))))
                self.y = max(0, min(self.screen_height - 80, int(settings.get("pet_y", self.y))))
                self.root.geometry(f"+{self.x}+{self.y}")
            except (TypeError, ValueError, tk.TclError):
                pass
        self.root.title(f"{self._pet_name()}教師秘書")
        self.attachments: list[Path] = []
        self.attachment_content_cache: dict[str, str] = {}
        self.draft: dict | None = None
        self.panel: tk.Toplevel | None = None
        self.panel_widgets: dict[str, object] = {}
        self.health_alerted = {"water": False, "move": False, "medicine": False}
        self.cheer_timer: str | None = None
        self.data_lock = threading.RLock()
        self.bridge_server: LocalBridgeServer | None = None
        self.bridge_thread: threading.Thread | None = None
        self.bridge_error = ""

        # 秘書面板改由瀏覽器呈現；右鍵與雙擊桌寵都直接開啟（不再使用 Tk 選單）。
        for widget in (self.root, self.pet_label):
            widget.bind("<Double-Button-1>", lambda _e: self.open_secretary_page())
        self.root.after(4000, self._visibility_heartbeat)
        self.root.after(1000, self._health_tick)
        self.root.after(1000, self._health_animation_tick)
        if settings.get("opening_brief", True):
            self.root.after(1200, self._opening_brief)
        self._schedule_cheer()
        self._start_local_bridge()
        if "--startup-launch" in sys.argv and settings.get("startup_delay_seconds", 0):
            self.root.withdraw()
            self.root.after(int(settings["startup_delay_seconds"]) * 1000, self.root.deiconify)
        if not settings.get("enabled", True):
            self.root.withdraw()
            self.root.after(100, self.root.destroy)

    def _pet_name(self) -> str:
        return str(self.data.get("settings", {}).get("pet_name") or "小綿助")

    def _in_quiet_hours(self) -> bool:
        settings = self.data.get("settings", {})
        start = str(settings.get("quiet_start") or "18:00")
        end = str(settings.get("quiet_end") or "07:30")
        current = datetime.now().strftime("%H:%M")
        if start == end:
            return False
        return start <= current < end if start < end else current >= start or current < end

    def _load_data(self) -> dict:
        DATA_ROOT.mkdir(parents=True, exist_ok=True)
        ATTACHMENT_ROOT.mkdir(parents=True, exist_ok=True)
        try:
            loaded = json.loads(DATA_FILE.read_text(encoding="utf-8"))
            if not isinstance(loaded, dict):
                raise ValueError("invalid data")
        except (OSError, ValueError, json.JSONDecodeError):
            loaded = default_data()
        base = default_data()
        base.update(loaded)
        if not isinstance(base.get("tasks"), list):
            base["tasks"] = []
        if not isinstance(base.get("notes"), list):
            base["notes"] = []
        if not isinstance(base.get("line_inbox"), list):
            base["line_inbox"] = []
        settings = default_data()["settings"]
        settings.update(base.get("settings") or {})
        settings.update(load_external_settings())
        settings["pet_name"] = str(settings.get("pet_name") or "小綿助").strip()[:12] or "小綿助"
        try:
            settings["cheer_interval_minutes"] = min(120, max(0, int(settings.get("cheer_interval_minutes", 5))))
        except (TypeError, ValueError):
            settings["cheer_interval_minutes"] = 5
        try:
            settings["startup_delay_seconds"] = min(300, max(0, int(settings.get("startup_delay_seconds", 30))))
        except (TypeError, ValueError):
            settings["startup_delay_seconds"] = 30
        for field in ("enabled", "autostart", "health_reminders", "opening_brief", "remember_position"):
            settings[field] = bool(settings.get(field, default_data()["settings"][field]))
        settings["bridge_enabled"] = bool(settings.get("bridge_enabled", True))
        settings["bridge_last_sync"] = str(settings.get("bridge_last_sync") or "").strip()
        allowed_origin = str(settings.get("allowed_web_origin") or "").strip().rstrip("/")
        settings["allowed_web_origin"] = allowed_origin if allowed_origin.startswith("https://") else ""
        for legacy_key in ("gas_sync_enabled", "gas_url", "gas_sync_key", "gas_last_sync"):
            settings.pop(legacy_key, None)
        settings["device_id"] = str(settings.get("device_id") or uuid.uuid4()).strip()
        for field, fallback in (("quiet_start", "18:00"), ("quiet_end", "07:30")):
            value = str(settings.get(field) or fallback)
            settings[field] = value if re.fullmatch(r"(?:[01]\d|2[0-3]):[0-5]\d", value) else fallback
        base["settings"] = settings
        health = default_data()["health"]
        health.update(base.get("health") or {})
        # 相容舊版單一服藥時間，升級成可多筆時間的格式。
        old_medicine_time = str(health.get("medicine_time") or "").strip()
        medicine_times = health.get("medicine_times")
        if not isinstance(medicine_times, list):
            medicine_times = []
        medicine_times = sorted({str(value).strip() for value in medicine_times if str(value).strip()})
        if old_medicine_time and old_medicine_time not in medicine_times:
            medicine_times.append(old_medicine_time)
            medicine_times.sort()
        health["medicine_times"] = medicine_times
        health["medicine_time"] = medicine_times[0] if medicine_times else ""
        try:
            health["move_interval"] = max(5, int(health.get("move_interval", 60)))
        except (TypeError, ValueError):
            health["move_interval"] = 60
        try:
            health["water_interval"] = max(5, int(health.get("water_interval", 60)))
        except (TypeError, ValueError):
            health["water_interval"] = 60
        for field in ("medicine_done_times", "medicine_alerted_times"):
            if not isinstance(health.get(field), list):
                health[field] = []
        if health["date"] != date.today().isoformat():
            health["date"] = date.today().isoformat()
            health["water_count"] = 0
            health["medicine_done_times"] = []
            health["medicine_alerted_times"] = []
        base["health"] = health
        return base

    def _save_data(self) -> None:
        with getattr(self, "data_lock", threading.RLock()):
            DATA_ROOT.mkdir(parents=True, exist_ok=True)
            temporary = DATA_FILE.with_suffix(".tmp")
            temporary.write_text(json.dumps(self.data, ensure_ascii=False, indent=2), encoding="utf-8")
            temporary.replace(DATA_FILE)

    def _finish_drag(self, _event: tk.Event) -> None:
        moved = self.drag_moved
        self.drag_origin = None
        self.play("idle")
        if moved and self.data.get("settings", {}).get("remember_position", True):
            self.data["settings"]["pet_x"] = int(self.x)
            self.data["settings"]["pet_y"] = int(self.y)
            self._save_data()
        if not moved:
            # 左鍵不再建立任何視窗，只播放一句打氣話。
            self._show_random_cheer()
        self.root.after(4200, self.start_walking)

    def _show_random_cheer(self) -> None:
        message = random.choice(self.CHEER_MESSAGES)
        self.play("success", 2, "idle")
        self.show_bubble(message, 4200)

    def _schedule_cheer(self) -> None:
        if self.cheer_timer:
            try:
                self.root.after_cancel(self.cheer_timer)
            except tk.TclError:
                pass
        minutes = min(120, max(0, int(self.data.get("settings", {}).get("cheer_interval_minutes", 5))))
        if minutes <= 0:
            self.cheer_timer = None
            return
        self.cheer_timer = self.root.after(minutes * 60_000, self._cheer_tick)

    def _cheer_tick(self) -> None:
        self.cheer_timer = None
        if not self._in_quiet_hours():
            self._show_random_cheer()
        self._schedule_cheer()

    def toggle_secretary(self) -> None:
        self.open_secretary_page()

    def open_secretary(self, brief_mode: str | None = None) -> None:
        # 舊 Tk 面板在部分機器（高解析縮放）會觸發繪圖層崩潰；一律改開瀏覽器版秘書頁。
        del brief_mode
        self.open_secretary_page()

    def _legacy_open_secretary(self, brief_mode: str | None = None) -> None:
        if not self.panel or not self.panel.winfo_exists():
            try:
                self._build_panel()
            except Exception:  # noqa: BLE001 - 建構失敗必須留下完整紀錄，避免只看到半張面板
                self._report_callback_exception(*sys.exc_info())
                raise
        if brief_mode:
            self._render_brief(brief_mode)
        self._render_dashboard()
        self.panel.deiconify()
        self.panel.lift()
        self.panel.focus_force()
        self._audit_panel_layout()
        # 高解析縮放螢幕上，Canvas 內嵌框架偶爾不會完成第一次重繪（看起來一片空白）。
        # 開啟後強制滾動一往一返，逼 Canvas 重新繪製全部內容。
        self.panel.after(180, self._force_panel_repaint)

    def _force_panel_repaint(self) -> None:
        try:
            canvas = next(c for c in self.panel.winfo_children() if isinstance(c, tk.Canvas))
            canvas.update_idletasks()
            canvas.yview_scroll(1, "units")
            canvas.update_idletasks()
            canvas.yview_moveto(0.0)
            canvas.update_idletasks()
        except (StopIteration, tk.TclError):
            pass

    def _audit_panel_layout(self) -> None:
        """面板健檢：卡片數量異常時把幾何狀態寫進錯誤日誌，協助遠端診斷。"""
        try:
            self.panel.update_idletasks()
            columns = self.panel_widgets.get("panel_columns") or ()
            detail = ["[panel-audit] 開啟面板診斷報告（雙欄版）",
                      f"panel geometry={self.panel.geometry()} state={self.panel.state()}"]
            for name, column in zip(("left", "right"), columns):
                detail.append(f"{name} col size={column.winfo_width()}x{column.winfo_height()} mapped={column.winfo_ismapped()}")
                for kid in column.winfo_children():
                    detail.append(f"  {kid.__class__.__name__} mgr={kid.winfo_manager()!r} "
                                  f"y={kid.winfo_y()} h={kid.winfo_height()} reqh={kid.winfo_reqheight()} mapped={kid.winfo_ismapped()}")
            ERROR_LOG.parent.mkdir(parents=True, exist_ok=True)
            with ERROR_LOG.open("a", encoding="utf-8") as stream:
                stream.write("\n[" + datetime.now().isoformat(timespec="seconds") + "]\n" + "\n".join(detail) + "\n")
        except Exception:  # noqa: BLE001 - 健檢本身失敗不影響面板使用
            pass

    def _build_panel(self) -> None:
        c = self.COLORS
        panel = tk.Toplevel(self.root)
        self.panel = panel
        panel.title(f"{self._pet_name()}教師秘書")
        panel.configure(bg=c["bg"])
        # 不把秘書視窗設為 topmost：Windows 的 Combobox 下拉層可能因此被主視窗擋住。
        panel.attributes("-topmost", False)
        _s = self.dpi_scale
        default_width = min(int(1180 * _s), max(int(1000 * _s), self.screen_width - 120))
        default_height = min(int(944 * _s), max(int(700 * _s), self.screen_height - 60))
        panel.geometry(self._panel_geometry(default_width, default_height))
        panel.minsize(int(980 * _s), int(700 * _s))
        panel.protocol("WM_DELETE_WINDOW", panel.withdraw)

        style = ttk.Style(panel)
        try:
            style.theme_use("clam")
        except tk.TclError:
            pass
        style.configure(
            "Secretary.TButton",
            font=("Microsoft JhengHei", 9, "bold"),
            padding=(11, 7),
            foreground=c["strong"],
            background=c["surface"],
            bordercolor=c["line"],
            lightcolor=c["line"],
            darkcolor=c["line"],
        )
        style.map("Secretary.TButton", background=[("active", c["soft"]), ("pressed", c["line"])])
        style.configure(
            "SecretaryPrimary.TButton",
            font=("Microsoft JhengHei", 9, "bold"),
            padding=(12, 8),
            foreground="white",
            background=c["primary"],
            bordercolor=c["primary"],
        )
        style.map("SecretaryPrimary.TButton", background=[("active", c["primary_hover"]), ("pressed", c["strong"])])
        style.configure("Secretary.TCombobox", padding=5, fieldbackground=c["surface"], background=c["surface"])
        style.configure("Secretary.Horizontal.TProgressbar", troughcolor="#dfe8e4", background=c["primary"], bordercolor="#dfe8e4")

        header = tk.Frame(panel, bg=c["primary"], padx=22, pady=9)
        header.pack(fill="x")
        header_copy = tk.Frame(header, bg=c["primary"])
        header_copy.pack(side="left", fill="x", expand=True)
        tk.Label(header_copy, text="TEACHER DESK SECRETARY", bg=c["primary"], fg="#bcd7d0", font=("Segoe UI", 8, "bold")).pack(anchor="w")
        header_title = tk.Label(header_copy, text=f"教師秘書・{self._pet_name()}", bg=c["primary"], fg="white", font=("Microsoft JhengHei", 19, "bold"))
        header_title.pack(anchor="w", pady=(2, 0))
        mode_badge = tk.Label(header, text="本機模式", bg="#dcece7", fg=c["strong"], padx=10, pady=5, font=("Microsoft JhengHei", 8, "bold"))
        mode_badge.pack(side="right", anchor="n")
        self.panel_widgets["header_title"] = header_title
        self.panel_widgets["mode_badge"] = mode_badge

        footer = tk.Frame(panel, bg=c["surface"], padx=18, pady=6, highlightbackground=c["line"], highlightthickness=1)
        footer.pack(side="bottom", fill="x")
        tk.Label(footer, text="資料保存在這台電腦・可匯出記事給教師工作台", bg=c["surface"], fg=c["muted"], font=("Microsoft JhengHei", 8)).pack(side="left")
        ttk.Button(footer, text="關閉", style="Secretary.TButton", command=panel.withdraw).pack(side="right")

        # 雙欄一頁式版面：不使用 Canvas 捲動容器（部分高解析縮放環境下
        # Canvas 內嵌框架會出現「元件存在但畫面空白」的重繪問題）。
        body = tk.Frame(panel, bg=c["bg"], padx=16, pady=12)
        body.pack(fill="both", expand=True)
        body.columnconfigure(0, weight=1, uniform="col")
        body.columnconfigure(1, weight=1, uniform="col")
        body.rowconfigure(0, weight=1)
        left_col = tk.Frame(body, bg=c["bg"])
        left_col.grid(row=0, column=0, sticky="nsew", padx=(0, 8))
        right_col = tk.Frame(body, bg=c["bg"])
        right_col.grid(row=0, column=1, sticky="nsew", padx=(8, 0))
        self.panel_widgets["panel_columns"] = (left_col, right_col)

        brief = self._card(left_col, "☀ 今日簡報", c["soft"])
        brief_label = tk.Label(brief, bg=c["soft"], fg=c["ink"], justify="left", anchor="w", wraplength=int(500 * self.dpi_scale), font=("Microsoft JhengHei", 10), pady=3)
        brief_label.pack(fill="x")
        brief_actions = tk.Frame(brief, bg=c["soft"])
        brief_actions.pack(fill="x", pady=(5, 0))
        ttk.Button(brief_actions, text="整理今天", style="Secretary.TButton", command=lambda: self._render_brief("today")).pack(side="left")
        ttk.Button(brief_actions, text="準備明天", style="Secretary.TButton", command=lambda: self._render_brief("tomorrow")).pack(side="left", padx=6)
        self.panel_widgets["brief"] = brief_label

        empty_state = tk.Frame(left_col, bg=c["surface"], padx=18, pady=14, highlightbackground=c["line"], highlightthickness=1)
        empty_state.pack(fill="x", pady=(6, 0))
        tk.Label(empty_state, text="目前還沒有任務", bg=c["surface"], fg=c["strong"], font=("Microsoft JhengHei", 12, "bold")).pack(anchor="w")
        tk.Label(
            empty_state,
            text="可以先在教師工作台新增任務，或使用下方「快速交代」建立本機記事；今日、逾期、待追蹤與健康管理仍可正常使用。",
            bg=c["surface"], fg=c["muted"], justify="left", anchor="w", wraplength=int(500 * self.dpi_scale),
            font=("Microsoft JhengHei", 9),
        ).pack(fill="x", pady=(5, 0))
        self.panel_widgets["empty_state"] = empty_state

        task_grid = tk.Frame(left_col, bg=c["bg"])
        task_grid.pack(fill="x", pady=6)
        self.panel_widgets["task_grid"] = task_grid
        task_grid.columnconfigure(0, weight=1)
        task_grid.columnconfigure(1, weight=1)
        today_card = self._card(task_grid, "📌 今日重要", c["surface"], pack=False)
        today_card.grid(row=0, column=0, sticky="nsew", padx=(0, 5))
        overdue_card = self._card(task_grid, "⚠ 已逾期", c["danger_soft"], pack=False)
        overdue_card.grid(row=0, column=1, sticky="nsew", padx=(5, 0))
        self.panel_widgets["today"] = self._listbox(today_card)
        self.panel_widgets["overdue"] = self._listbox(overdue_card)

        tracking = self._card(left_col, "📮 待追蹤｜等待回覆、尚未收齊", c["surface"])
        tracking.pack(fill="x", pady=(0, 6))
        self.panel_widgets["tracking"] = self._listbox(tracking, height=3)

        health = self._card(left_col, "🌿 健康管理", c["soft"])
        health.pack(fill="x", pady=(0, 6))
        self._build_health_rows(health)

        personalization = self._card(right_col, "⚙ 個人化設定", c["surface"])
        personalization.pack(fill="x", pady=(0, 6))
        settings_row = tk.Frame(personalization, bg=c["surface"])
        settings_row.pack(fill="x")
        settings_row.columnconfigure(1, weight=1)
        settings_row.columnconfigure(3, weight=1)
        tk.Label(settings_row, text="名稱", bg=c["surface"], fg=c["ink"], font=("Microsoft JhengHei", 9, "bold")).grid(row=0, column=0, sticky="w")
        pet_name_var = tk.StringVar(value=self._pet_name())
        ttk.Entry(settings_row, textvariable=pet_name_var, width=14).grid(row=0, column=1, sticky="ew", padx=(7, 18))
        tk.Label(settings_row, text="打氣間隔", bg=c["surface"], fg=c["ink"], font=("Microsoft JhengHei", 9, "bold")).grid(row=0, column=2, sticky="w")
        cheer_interval_var = tk.StringVar(value=str(self.data["settings"].get("cheer_interval_minutes", 5)))
        ttk.Combobox(settings_row, textvariable=cheer_interval_var, values=["1", "3", "5", "10", "15", "30", "60"], state="normal", width=6, style="Secretary.TCombobox").grid(row=0, column=3, sticky="ew", padx=(7, 5))
        tk.Label(settings_row, text="分鐘", bg=c["surface"], fg=c["muted"], font=("Microsoft JhengHei", 8)).grid(row=0, column=4, sticky="w")
        ttk.Button(settings_row, text="儲存設定", style="Secretary.TButton", command=self._save_personalization).grid(row=0, column=5, padx=(12, 0))
        tk.Label(personalization, text="左鍵點桌寵會顯示打氣話；右鍵可開啟這個秘書首頁。", bg=c["surface"], fg=c["muted"], font=("Microsoft JhengHei", 8)).pack(anchor="w", pady=(8, 0))
        self.panel_widgets["pet_name_var"] = pet_name_var
        self.panel_widgets["cheer_interval_var"] = cheer_interval_var

        bridge = self._card(right_col, "🔗 教師工作台本機連線", c["surface"])
        bridge.pack(fill="x", pady=(0, 6))
        tk.Label(
            bridge,
            text="小綿助開啟時，Netlify 或本機教師工作台會透過這台電腦的 127.0.0.1 自動交換任務與記事；資料不會經過第三方伺服器。",
            bg=c["surface"], fg=c["muted"], justify="left", anchor="w", wraplength=int(500 * self.dpi_scale),
            font=("Microsoft JhengHei", 8),
        ).pack(fill="x", pady=(0, 8))
        bridge_status = tk.Label(bridge, bg=c["surface"], fg=c["muted"], justify="left", anchor="w", font=("Microsoft JhengHei", 8, "bold"))
        bridge_status.pack(fill="x")
        self.panel_widgets["bridge_status"] = bridge_status
        self._update_bridge_status()

        quick = self._card(right_col, "💬 快速交代", c["surface"])
        quick.pack(fill="x")
        tk.Label(quick, text="輸入文字，或加入圖片、PDF、Word、Excel、PowerPoint，讓秘書讀取後整理。", bg=c["surface"], fg=c["muted"], font=("Microsoft JhengHei", 8)).pack(anchor="w", pady=(0, 7))
        handoff = tk.Text(quick, height=3, wrap="word", relief="flat", borderwidth=0, highlightthickness=1, highlightbackground=c["line"], highlightcolor=c["primary"], padx=9, pady=8, font=("Microsoft JhengHei", 10), undo=True)
        handoff.pack(fill="x")
        handoff.bind("<Control-v>", self._paste_attachment, add="+")
        self.panel_widgets["handoff"] = handoff
        attachment_label = tk.Label(quick, text="尚未加入附件", bg=c["surface"], fg=c["muted"], anchor="w", justify="left", wraplength=int(500 * self.dpi_scale), font=("Microsoft JhengHei", 8))
        attachment_label.pack(fill="x", pady=5)
        self.panel_widgets["attachments"] = attachment_label
        actions = tk.Frame(quick, bg=c["surface"])
        actions.pack(fill="x")
        ttk.Button(actions, text="📎 加入附件", style="Secretary.TButton", command=self._choose_attachments).pack(side="left")
        ttk.Button(actions, text="清除", style="Secretary.TButton", command=self._clear_attachments).pack(side="left", padx=5)
        ttk.Button(actions, text="🎙 語音", style="Secretary.TButton", command=self._voice_handoff).pack(side="left")
        primary_actions = tk.Frame(quick, bg=c["surface"])
        primary_actions.pack(fill="x", pady=(8, 0))
        organize_button = ttk.Button(primary_actions, text=f"交給{self._pet_name()}整理", style="SecretaryPrimary.TButton", command=self._analyze_handoff)
        organize_button.pack(side="right")
        ttk.Button(primary_actions, text="先存為記事", style="Secretary.TButton", command=self._save_note).pack(side="right", padx=7)
        self.panel_widgets["organize_button"] = organize_button

        line_card = self._card(right_col, "📱 LINE 待整理", c["surface"])
        line_card.pack(fill="x", pady=(6, 0))
        tk.Label(line_card, text="從 LINE 傳給小幫手的訊息；請到教師工作台的「LINE 收件匣」確認後建立。", bg=c["surface"], fg=c["muted"], font=("Microsoft JhengHei", 8)).pack(anchor="w", pady=(0, 4))
        line_list = self._listbox(line_card, height=3)
        self.panel_widgets["line_inbox"] = line_list

        notes = self._card(right_col, "📝 本機記事", c["surface"])
        notes.pack(fill="x", pady=(6, 0))
        tk.Label(notes, text="記事保存在這台電腦；教師工作台開啟時會透過本機連線自動帶入。", bg=c["surface"], fg=c["muted"], font=("Microsoft JhengHei", 8)).pack(anchor="w", pady=(0, 4))
        notes_list = self._listbox(notes, height=3)
        note_actions = tk.Frame(notes, bg=c["surface"])
        note_actions.pack(fill="x", pady=(5, 0))
        ttk.Button(note_actions, text="開啟記事資料夾", style="Secretary.TButton", command=self._open_data_folder).pack(side="right")
        ttk.Button(note_actions, text="匯出給網頁", style="Secretary.TButton", command=self._export_notes_for_web).pack(side="right", padx=5)
        self.panel_widgets["notes"] = notes_list

        draft = tk.Frame(right_col, bg=c["gold_soft"], padx=12, pady=10, highlightbackground="#d9c58e", highlightthickness=1)
        draft_label = tk.Label(draft, bg=c["gold_soft"], fg=c["ink"], justify="left", anchor="w", wraplength=int(500 * self.dpi_scale), font=("Microsoft JhengHei", 10, "bold"))
        draft_label.pack(fill="x")
        draft_actions = tk.Frame(draft, bg=c["gold_soft"])
        draft_actions.pack(fill="x", pady=(8, 0))
        ttk.Button(draft_actions, text="確認建立本機任務", style="Secretary.TButton", command=self._confirm_draft).pack(side="left")
        ttk.Button(draft_actions, text="返回修改", style="Secretary.TButton", command=draft.pack_forget).pack(side="left", padx=6)
        self.panel_widgets["draft_frame"] = draft
        self.panel_widgets["draft_label"] = draft_label

        panel.bind("<Escape>", lambda _e: panel.withdraw())
        self._render_brief("tomorrow" if datetime.now().hour >= 15 else "today")
        self._render_health()
        self._render_line_inbox()

    def _panel_geometry(self, width: int, height: int) -> str:
        max_x = max(10, self.screen_width - width - 14)
        max_y = max(10, self.screen_height - height - 70)
        x = min(max_x, max(10, self.x - width + 130))
        y = min(max_y, max(10, self.y - height + 180))
        return f"{width}x{height}+{x}+{y}"

    def _card(self, parent: tk.Misc, title: str, bg: str, pack: bool = True) -> tk.Frame:
        frame = tk.Frame(parent, bg=bg, padx=14, pady=12, highlightbackground=self.COLORS["line"], highlightthickness=1)
        if pack:
            frame.pack(fill="x")
        tk.Label(frame, text=title, bg=bg, fg=self.COLORS["strong"], font=("Microsoft JhengHei", 11, "bold")).pack(anchor="w", pady=(0, 7))
        return frame

    def _listbox(self, parent: tk.Misc, height: int = 4) -> tk.Listbox:
        bg = str(parent.cget("bg"))
        box = tk.Listbox(
            parent,
            height=height,
            borderwidth=0,
            highlightthickness=0,
            activestyle="none",
            bg=bg,
            fg=self.COLORS["ink"],
            selectbackground=self.COLORS["soft"],
            selectforeground=self.COLORS["strong"],
            font=("Microsoft JhengHei", 9),
        )
        # Keep lists at their requested row height. Expanding them inside the
        # scrolling canvas can consume the remaining window and look like a
        # large blank panel on high-DPI Windows displays.
        box.pack(fill="x", expand=False)
        return box

    def _build_health_rows(self, parent: tk.Frame) -> None:
        c = self.COLORS
        # 喝水：用水滴圖示呈現今日累計，按鈕只負責增加一杯。
        water_row = tk.Frame(parent, bg=c["soft"], pady=6)
        water_row.pack(fill="x")
        tk.Label(water_row, text="💧 喝水", bg=c["soft"], fg=c["ink"], width=12, anchor="w", font=("Microsoft JhengHei", 9, "bold")).pack(side="left")
        water_info = tk.Frame(water_row, bg=c["soft"])
        water_info.pack(side="left", fill="x", expand=True)
        water_visual = tk.Label(water_info, bg=c["soft"], fg=c["primary"], anchor="w", font=("Segoe UI Emoji", 11))
        water_visual.pack(anchor="w")
        water_label = tk.Label(water_info, bg=c["soft"], fg=c["muted"], anchor="w", font=("Microsoft JhengHei", 8))
        water_label.pack(anchor="w")
        ttk.Button(water_row, text="喝一杯", style="Secretary.TButton", command=lambda: self._record_health("water")).pack(side="right")
        self.panel_widgets["health_water"] = water_label
        self.panel_widgets["health_water_visual"] = water_visual

        # 起身活動：進度條每秒更新，顯示距離下一次提醒的累計時間。
        move_row = tk.Frame(parent, bg=c["soft"], pady=7)
        move_row.pack(fill="x")
        tk.Label(move_row, text="🚶 起身活動", bg=c["soft"], fg=c["ink"], width=12, anchor="nw", font=("Microsoft JhengHei", 9, "bold")).pack(side="left")
        move_info = tk.Frame(move_row, bg=c["soft"])
        move_info.pack(side="left", fill="x", expand=True, padx=(0, 6))
        move_timer = tk.Label(move_info, bg=c["soft"], fg=c["muted"], anchor="w", font=("Microsoft JhengHei", 8))
        move_timer.pack(fill="x")
        move_progress = ttk.Progressbar(move_info, orient="horizontal", mode="determinate", maximum=100, length=170, style="Secretary.Horizontal.TProgressbar")
        move_progress.pack(fill="x", pady=(2, 0))
        move_controls = tk.Frame(move_row, bg=c["soft"])
        move_controls.pack(side="right")
        interval_var = tk.StringVar(value=str(self.data["health"].get("move_interval", 60)))
        interval_box = ttk.Combobox(move_controls, textvariable=interval_var, values=["15", "30", "45", "60", "90", "120"], state="readonly", width=5, style="Secretary.TCombobox")
        interval_box.pack(side="left", padx=(0, 4))
        interval_box.bind("<<ComboboxSelected>>", lambda _e: self._set_move_interval(interval_var.get()))
        tk.Label(move_controls, text="分", bg=c["soft"], fg=c["muted"], font=("Microsoft JhengHei", 8)).pack(side="left")
        ttk.Button(move_controls, text="活動完成", style="Secretary.TButton", command=lambda: self._record_health("move")).pack(side="left", padx=(5, 0))
        self.panel_widgets["health_move"] = move_timer
        self.panel_widgets["health_move_timer"] = move_timer
        self.panel_widgets["health_move_progress"] = move_progress

        # 服藥：時間只能從下拉選單選擇，但可新增多筆提醒。
        medicine_row = tk.Frame(parent, bg=c["soft"], pady=7)
        medicine_row.pack(fill="x")
        tk.Label(medicine_row, text="💊 服藥", bg=c["soft"], fg=c["ink"], width=12, anchor="nw", font=("Microsoft JhengHei", 9, "bold")).pack(side="left")
        medicine_body = tk.Frame(medicine_row, bg=c["soft"])
        medicine_body.pack(side="left", fill="x", expand=True)
        medicine_control = tk.Frame(medicine_body, bg=c["soft"])
        medicine_control.pack(fill="x")
        med_var = tk.StringVar(value="")
        med_values = [f"{hour:02d}:{minute:02d}" for hour in range(24) for minute in (0, 30)]
        med_box = ttk.Combobox(medicine_control, textvariable=med_var, values=med_values, state="readonly", width=8, style="Secretary.TCombobox", height=10)
        med_box.pack(side="left")
        ttk.Button(medicine_control, text="＋新增時間", style="Secretary.TButton", command=self._add_medicine_time).pack(side="left", padx=5)
        tk.Label(medicine_body, text="可設定多個提醒時間；每一筆都由小綿助提醒。", bg=c["soft"], fg=c["muted"], anchor="w", font=("Microsoft JhengHei", 8)).pack(fill="x")
        medicine_list = tk.Frame(medicine_body, bg=c["soft"])
        medicine_list.pack(fill="x", pady=(3, 0))
        self.panel_widgets["health_medicine"] = medicine_list
        self.panel_widgets["medicine_combo"] = med_box
        self.panel_widgets["medicine_var"] = med_var
        self.panel_widgets["medicine_list"] = medicine_list
        self._render_medicine_list()

    def _render_dashboard(self) -> None:
        if not self.panel_widgets:
            return
        today = date.today().isoformat()
        active = [task for task in self.data["tasks"] if task.get("status") != "已完成"]
        due_today = [task for task in active if task.get("due_date") == today]
        overdue = [task for task in active if task.get("due_date") and task["due_date"] < today]
        tracking = [task for task in active if task.get("waiting_for") or task.get("status") == "等待回覆"]
        empty_state = self.panel_widgets.get("empty_state")
        if isinstance(empty_state, tk.Frame):
            if self.data["tasks"]:
                empty_state.pack_forget()
            elif not empty_state.winfo_manager():
                empty_state.pack(fill="x", pady=(6, 0), before=self.panel_widgets["task_grid"])
        self._fill_task_list("today", due_today, "今天沒有期限任務")
        self._fill_task_list("overdue", overdue, "目前沒有逾期任務")
        self._fill_task_list("tracking", tracking, "目前沒有等待回覆的事項")
        self._render_health()
        self._render_notes()

    def _save_personalization(self) -> None:
        name_var: tk.StringVar = self.panel_widgets["pet_name_var"]  # type: ignore[assignment]
        interval_var: tk.StringVar = self.panel_widgets["cheer_interval_var"]  # type: ignore[assignment]
        name = name_var.get().strip()[:12] or "小綿助"
        try:
            interval = min(120, max(1, int(interval_var.get())))
        except (TypeError, ValueError):
            interval = 5
        self.data["settings"]["pet_name"] = name
        self.data["settings"]["cheer_interval_minutes"] = interval
        self._save_data()
        self.root.title(f"{name}教師秘書")
        if self.panel and self.panel.winfo_exists():
            self.panel.title(f"{name}教師秘書")
        header: tk.Label = self.panel_widgets["header_title"]  # type: ignore[assignment]
        header.configure(text=f"教師秘書・{name}")
        organize_button: ttk.Button = self.panel_widgets["organize_button"]  # type: ignore[assignment]
        organize_button.configure(text=f"交給{name}整理")
        self.menu.entryconfigure(0, label=f"開啟{name}秘書")
        self._schedule_cheer()
        self.play("success", 2, "idle", f"好！以後可以叫我{name}。")

    # ------------------------------------------------------------------
    # 瀏覽器版秘書面板：資料與動作由本機橋接供應，畫面交給瀏覽器渲染。
    # （部分機器的 Tk 文字繪製不穩定；瀏覽器在任何顯示環境都可靠。）
    # ------------------------------------------------------------------

    def _show_menu(self, _event: tk.Event) -> str:
        # Tk 右鍵選單與面板同屬會在部分機器崩潰／消失的繪圖管線；
        # 右鍵改為直接開啟瀏覽器版秘書（暫停、關閉等功能都在秘書頁上）。
        self.open_secretary_page()
        return "break"

    def _visibility_heartbeat(self) -> None:
        """部分顯示卡會無聲丟棄透明色鍵視窗的畫面；定期重新亮相讓桌寵自我復活。"""
        try:
            if self.data.get("settings", {}).get("enabled", True):
                self.root.deiconify()
                self.root.attributes("-topmost", True)
                self.root.lift()
        except tk.TclError:
            return
        self.root.after(4000, self._visibility_heartbeat)

    def open_secretary_page(self) -> None:
        if self.bridge_error:
            self.show_bubble("本機服務未啟動，秘書頁開不了；請重新啟動小綿助。", 3200)
            return
        import webbrowser

        webbrowser.open(f"http://{BRIDGE_HOST}:{BRIDGE_PORT}/panel")
        self.play("success", 2, "idle", "秘書頁開好了！")

    def _panel_payload(self) -> dict:
        today = date.today().isoformat()
        with self.data_lock:
            active = [task for task in self.data.get("tasks", []) if task.get("status") != "已完成"]
            def pick(tasks):
                return [{"title": str(t.get("title") or "未命名")[:80], "due_time": str(t.get("due_time") or "")} for t in tasks[:8]]
            due_today = pick([t for t in active if t.get("due_date") == today])
            overdue = pick([t for t in active if t.get("due_date") and t["due_date"] < today])
            tracking = pick([t for t in active if t.get("waiting_for") or t.get("status") == "等待回覆"])
            notes = [{"text": str(n.get("text") or "")[:200], "time": str(n.get("created_at") or "")[5:16].replace("T", " ")}
                     for n in self.data.get("notes", [])[:10] if n.get("text")]
            line_items = [{"title": str(i.get("title") or "")[:80],
                           "kind": "任務" if i.get("kind") == "task" or i.get("type") == "task" else "記事",
                           "tag": str(i.get("tag") or ""), "medium": str(i.get("medium") or "")}
                          for i in self.data.get("line_inbox", []) if isinstance(i, dict)]
            health = self.data.get("health", {})
            try:
                last_move = datetime.fromisoformat(str(health.get("last_move")))
                elapsed = int((datetime.now() - last_move).total_seconds() // 60)
            except (TypeError, ValueError):
                elapsed = 0
            payload = {
                "ok": True,
                "petName": self._pet_name(),
                "generatedAt": datetime.now().strftime("%H:%M:%S"),
                "dueToday": due_today, "overdue": overdue, "tracking": tracking,
                "notes": notes, "lineInbox": line_items,
                "health": {
                    "water": int(health.get("water_count") or 0),
                    "moveElapsed": max(0, elapsed),
                    "moveInterval": int(health.get("move_interval") or 60),
                    "waterInterval": int(health.get("water_interval") or 60),
                    "moveDone": health.get("move_done_date") == today,
                    "medicineTimes": list(health.get("medicine_times") or []),
                    "medicineDone": list(health.get("medicine_done_times") or []),
                },
            }
        return payload

    def _panel_apply_action(self, payload: dict) -> dict:
        action = str(payload.get("action") or "")
        now = now_iso()
        today = date.today().isoformat()
        bubble = ""
        with self.data_lock:
            health = self.data.setdefault("health", {})
            if action == "water":
                health["water_count"] = int(health.get("water_count") or 0) + 1
                health["last_water"] = now
                bubble = f"喝水第 {health['water_count']} 杯，讚！"
            elif action == "move_done":
                health["last_move"] = now
                health["move_done_date"] = today
                bubble = "起身活動完成，繼續加油！"
            elif action == "medicine_done":
                slot = str(payload.get("time") or "")
                done = health.setdefault("medicine_done_times", [])
                if slot and slot not in done:
                    done.append(slot)
                bubble = f"{slot} 的藥記錄好了。"
            elif action == "set_move_interval":
                try:
                    minutes = min(240, max(5, int(payload.get("minutes"))))
                except (TypeError, ValueError):
                    return {"ok": False, "error": "間隔必須是 5～240 的分鐘數"}
                health["move_interval"] = minutes
                bubble = f"起身提醒改為每 {minutes} 分鐘。"
            elif action == "set_water_interval":
                try:
                    minutes = min(240, max(5, int(payload.get("minutes"))))
                except (TypeError, ValueError):
                    return {"ok": False, "error": "間隔必須是 5～240 的分鐘數"}
                health["water_interval"] = minutes
                bubble = f"喝水提醒改為每 {minutes} 分鐘。"
            elif action == "add_medicine_time":
                slot = str(payload.get("time") or "").strip()
                if not re.fullmatch(r"([01]\d|2[0-3]):[0-5]\d", slot):
                    return {"ok": False, "error": "服藥時間格式須為 HH:MM"}
                times = health.setdefault("medicine_times", [])
                if slot not in times:
                    times.append(slot)
                    times.sort()
                if len(times) > 12:
                    times[:] = times[:12]
                health["medicine_time"] = times[0] if times else ""
                bubble = f"已加入 {slot} 的服藥提醒。"
            elif action == "remove_medicine_time":
                slot = str(payload.get("time") or "").strip()
                times = health.setdefault("medicine_times", [])
                if slot in times:
                    times.remove(slot)
                for field in ("medicine_done_times", "medicine_alerted_times"):
                    values = health.get(field)
                    if isinstance(values, list) and slot in values:
                        values.remove(slot)
                health["medicine_time"] = times[0] if times else ""
                bubble = f"已移除 {slot} 的服藥提醒。"
            elif action == "add_note":
                text = str(payload.get("text") or "").strip()[:3000]
                if not text:
                    return {"ok": False, "error": "記事內容不可空白"}
                self.data.setdefault("notes", []).insert(0, {
                    "id": str(uuid.uuid4()), "text": text, "attachments": [], "created_at": now,
                })
                bubble = "記事保存好了。"
            elif action == "toggle_pause":
                try:
                    self.root.after(0, self.toggle_pause)
                except tk.TclError:
                    pass
                return {"ok": True, "message": "已切換暫停／繼續"}
            elif action == "quit":
                try:
                    self.root.after(150, self.root.destroy)
                except tk.TclError:
                    pass
                return {"ok": True, "message": "小綿助即將關閉"}
            else:
                return {"ok": False, "error": f"未知動作：{action}"}
            self._save_data()
        try:
            self.root.after(0, lambda: self.show_bubble(bubble, 2600))
        except tk.TclError:
            pass
        return self._panel_payload()

    def _panel_html(self) -> str:
        return PANEL_PAGE_HTML.replace("__PET_NAME__", self._pet_name())

    def _start_local_bridge(self) -> None:
        if not self.data.get("settings", {}).get("bridge_enabled", True):
            return
        try:
            self.bridge_server = LocalBridgeServer((BRIDGE_HOST, BRIDGE_PORT), make_bridge_handler(self))
            self.bridge_thread = threading.Thread(target=self.bridge_server.serve_forever, name="xiaomianzhu-local-bridge", daemon=True)
            self.bridge_thread.start()
        except OSError as error:
            # 8767 被占用最常見的原因是「已有另一個小綿助在執行」。多個實例會疊在
            # 同一位置、面板互相干擾，看起來像壞掉；直接提示並結束，避免殭屍實例。
            self.bridge_error = str(error)
            try:
                messagebox.showwarning(
                    APP_NAME,
                    "偵測到本機連接埠 8767 已被使用。\n\n"
                    "最可能的原因：小綿助已經在執行中（請看工作列或桌面角落，"
                    "或用工作管理員結束多餘的 XiaoMianZhu／python 程序後再開一次）。\n\n"
                    "若是其他程式占用了 8767，關閉該程式後重新啟動小綿助即可。\n"
                    "本次啟動將結束，以免出現多隻小綿助互相干擾。",
                )
            except tk.TclError:
                pass
            self.root.destroy()
            sys.exit(0)

    def _update_bridge_status(self) -> None:
        badge = self.panel_widgets.get("mode_badge")
        if isinstance(badge, tk.Label):
            badge.configure(text="本機連線模式", bg="#dcece7")
        status = self.panel_widgets.get("bridge_status")
        if isinstance(status, tk.Label):
            if self.bridge_error:
                status.configure(text=f"本機連線未啟動：{self.bridge_error}", fg=self.COLORS["danger"])
                return
            last_sync = str(self.data.get("settings", {}).get("bridge_last_sync") or "").replace("T", " ")
            suffix = f"｜最後交換：{last_sync}" if last_sync else "｜等待教師工作台連線"
            status.configure(text=f"已在 http://{BRIDGE_HOST}:{BRIDGE_PORT} 啟動{suffix}", fg=self.COLORS["muted"])

    def _bridge_health(self) -> dict:
        return {
            "ok": True,
            "service": "xiaomianzhu-local-bridge",
            "version": 1,
            "petName": self._pet_name(),
            "lastSync": self.data.get("settings", {}).get("bridge_last_sync", ""),
        }

    def _bridge_allowed_origin(self) -> str:
        return str(self.data.get("settings", {}).get("allowed_web_origin") or "")

    @staticmethod
    def _web_task_to_local(item: dict) -> dict:
        cloud_id = str(item.get("taskId") or "").strip()
        return {
            "id": cloud_id,
            "cloud_id": cloud_id,
            "title": str(item.get("name") or "未命名任務")[:160],
            "kind": "行程" if item.get("dueTime") else "任務",
            "due_date": str(item.get("dueDate") or ""),
            "due_time": str(item.get("dueTime") or ""),
            "priority": str(item.get("priority") or "中"),
            "status": str(item.get("status") or "未開始"),
            "waiting_for": str(item.get("waitingFor") or ""),
            "source": "web",
            "created_at": str(item.get("createdAt") or now_iso()),
            "web_task": dict(item),
        }

    @staticmethod
    def _local_task_to_web(task: dict) -> dict:
        raw = task.get("web_task") if isinstance(task.get("web_task"), dict) else {}
        task_id = str(task.get("cloud_id") or raw.get("taskId") or f"PET-{task.get('id') or uuid.uuid4()}")
        try:
            sort_order = int(raw.get("sortOrder") or 999)
        except (TypeError, ValueError):
            sort_order = 999
        return {
            **raw,
            "taskId": task_id,
            "name": str(task.get("title") or raw.get("name") or "未命名任務")[:160],
            "category": str(raw.get("category") or "其他"),
            "status": str(task.get("status") or raw.get("status") or "未開始"),
            "priority": str(task.get("priority") or raw.get("priority") or "中"),
            "dueDate": str(task.get("due_date") or raw.get("dueDate") or ""),
            "dueTime": str(task.get("due_time") or raw.get("dueTime") or ""),
            "nextAction": str(raw.get("nextAction") or task.get("title") or ""),
            "waitingFor": str(task.get("waiting_for") or raw.get("waitingFor") or ""),
            "owner": str(raw.get("owner") or "桌面小綿助"),
            "boardDisplay": str(raw.get("boardDisplay") or "自動"),
            "sortOrder": sort_order,
            "source": "desktop" if task.get("source") != "web" else "web",
        }

    def _bridge_sync(self, payload: dict) -> dict:
        incoming_tasks = payload.get("tasks") if isinstance(payload.get("tasks"), list) else []
        incoming_notes = payload.get("notes") if isinstance(payload.get("notes"), list) else []
        incoming_line = payload.get("lineInbox") if isinstance(payload.get("lineInbox"), list) else None
        with self.data_lock:
            if incoming_line is not None:
                # 網頁是 LINE 收件匣的處理者；桌面端只保存唯讀鏡像供顯示與簡報。
                self.data["line_inbox"] = sanitize_line_inbox(incoming_line)
            local_tasks = self.data.get("tasks", [])
            by_cloud = {str(task.get("cloud_id") or ""): task for task in local_tasks if task.get("cloud_id")}
            for item in incoming_tasks[:5000]:
                if not isinstance(item, dict) or not item.get("taskId"):
                    continue
                converted = self._web_task_to_local(item)
                existing = by_cloud.get(converted["cloud_id"])
                if existing:
                    existing.update(converted)
                else:
                    local_tasks.append(converted)
                    by_cloud[converted["cloud_id"]] = converted

            local_notes = self.data.get("notes", [])
            note_keys = {f"{note.get('created_at', '')}|{note.get('text', '')}" for note in local_notes}
            for note in incoming_notes[:500]:
                if not isinstance(note, dict):
                    continue
                text = str(note.get("text") or "").strip()[:3000]
                created_at = str(note.get("createdAt") or note.get("created_at") or now_iso())
                key = f"{created_at}|{text}"
                if text and key not in note_keys:
                    attachments = note.get("attachments") if isinstance(note.get("attachments"), list) else []
                    local_notes.append({
                        "id": str(note.get("id") or uuid.uuid4()),
                        "text": text,
                        "attachments": [str(value)[:180] for value in attachments[:20]],
                        "created_at": created_at,
                        "source": "web",
                    })
                    note_keys.add(key)
            local_notes.sort(key=lambda note: str(note.get("created_at") or ""), reverse=True)
            self.data["notes"] = local_notes[:500]
            self.data["settings"]["bridge_last_sync"] = now_iso()
            self._save_data()
            response_tasks = [self._local_task_to_web(task) for task in self.data.get("tasks", [])[:5000]]
            response_notes = []
            for note in self.data.get("notes", []):
                if not note.get("text") and not note.get("attachment_text"):
                    continue
                note_id = str(note.get("id") or uuid.uuid4())
                note["id"] = note_id
                attachments = note.get("attachments") if isinstance(note.get("attachments"), list) else []
                response_notes.append({
                    "id": note_id,
                    "text": str(note.get("text") or note.get("attachment_text") or "")[:3000],
                    "attachments": [str(value)[:180] for value in attachments[:20]],
                    "createdAt": str(note.get("created_at") or now_iso()),
                    "source": "desktop" if note.get("source") != "web" else "web",
                })
        try:
            self.root.after(0, self._bridge_ui_refresh)
        except tk.TclError:
            pass
        return {"ok": True, "tasks": response_tasks, "notes": response_notes, "syncedAt": now_iso()}

    def _bridge_ui_refresh(self) -> None:
        self._render_dashboard()
        self._render_notes()
        self._render_line_inbox()
        self._update_bridge_status()

    def _render_line_inbox(self) -> None:
        box = self.panel_widgets.get("line_inbox")
        if not isinstance(box, tk.Listbox):
            return
        box.delete(0, "end")
        items = [item for item in self.data.get("line_inbox", []) if isinstance(item, dict)]
        if not items:
            box.insert("end", "目前沒有 LINE 待整理的訊息")
            return
        for item in items[:6]:
            kind = "任務" if item.get("type") == "task" else "記事"
            medium = "🎤" if item.get("medium") == "voice" else "📷" if item.get("medium") == "photo" else ""
            tag = f"（{item.get('tag')}）" if item.get("tag") else ""
            box.insert("end", f"{kind}{medium}｜{item.get('title', '')}{tag}")

    def _render_notes(self) -> None:
        box = self.panel_widgets.get("notes")
        if not isinstance(box, tk.Listbox):
            return
        box.delete(0, "end")
        notes = self.data.get("notes", [])
        if not notes:
            box.insert("end", "目前沒有本機記事")
            return
        for note in notes[:8]:
            timestamp = str(note.get("created_at", ""))[5:16].replace("T", " ")
            text = re.sub(r"\s+", " ", str(note.get("text") or "附件記事")).strip()
            box.insert("end", f"{timestamp}｜{text[:58]}")

    def _open_data_folder(self) -> None:
        DATA_ROOT.mkdir(parents=True, exist_ok=True)
        try:
            os.startfile(DATA_ROOT)  # type: ignore[attr-defined]
        except (AttributeError, OSError):
            messagebox.showinfo(APP_NAME, str(DATA_ROOT), parent=self.panel)

    def _export_notes_for_web(self) -> None:
        notes = self.data.get("notes", [])
        if not notes:
            messagebox.showinfo(APP_NAME, "目前沒有可以匯出的記事。", parent=self.panel)
            return
        target = filedialog.asksaveasfilename(
            parent=self.panel,
            title="匯出小綿助記事給教師工作台",
            defaultextension=".json",
            filetypes=[("JSON 檔案", "*.json")],
            initialfile=f"小綿助記事-{date.today().isoformat()}.json",
        )
        if not target:
            return
        safe_notes = []
        for note in notes[:500]:
            text = str(note.get("text") or note.get("attachment_text") or "").strip()[:3000]
            if not text:
                continue
            safe_notes.append({
                "id": str(note.get("id") or uuid.uuid4()),
                "text": text,
                "attachments": [str(name)[:180] for name in note.get("attachments", [])[:20]],
                "createdAt": str(note.get("created_at") or now_iso()),
                "source": "desktop",
            })
        payload = {
            "kind": "teacher-dashboard-secretary-notes",
            "version": 1,
            "exportedAt": now_iso(),
            "notes": safe_notes,
        }
        Path(target).write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        messagebox.showinfo(
            APP_NAME,
            f"已匯出 {len(safe_notes)} 則記事。\n請到教師工作台的「小綿助記事」按「匯入桌面記事」。",
            parent=self.panel,
        )

    def _fill_task_list(self, key: str, tasks: list[dict], empty: str) -> None:
        box: tk.Listbox = self.panel_widgets[key]  # type: ignore[assignment]
        box.delete(0, "end")
        if not tasks:
            box.insert("end", empty)
            return
        for task in tasks[:5]:
            suffix = f"｜{task.get('due_time')}" if task.get("due_time") else ""
            box.insert("end", f"{task.get('title', '未命名')}{suffix}")

    def _render_brief(self, mode: str) -> None:
        if "brief" not in self.panel_widgets:
            return
        target = date.today() + (timedelta(days=1) if mode == "tomorrow" else timedelta())
        active = [task for task in self.data["tasks"] if task.get("status") != "已完成"]
        due = [task for task in active if task.get("due_date") == target.isoformat()]
        overdue = [task for task in active if task.get("due_date") and task["due_date"] < date.today().isoformat()] if mode == "today" else []
        tracking = [task for task in active if task.get("waiting_for") or task.get("status") == "等待回覆"]
        line_pending = len([item for item in self.data.get("line_inbox", []) if isinstance(item, dict)])
        if mode == "tomorrow":
            text = f"明天有 {len(due)} 件期限任務。\n" + (f"建議先準備：{due[0]['title']}" if due else "目前沒有需要提前準備的期限事項。")
        else:
            text = f"今天有 {len(due)} 件期限任務、{len(overdue)} 件逾期、{len(tracking)} 件等待回覆。\n"
            text += f"建議優先處理：{overdue[0]['title']}" if overdue else f"今天先完成：{due[0]['title']}" if due else "今天的期限工作已整理完成。"
        if line_pending:
            text += f"\n另有 {line_pending} 則 LINE 訊息待到工作台整理。"
        label: tk.Label = self.panel_widgets["brief"]  # type: ignore[assignment]
        label.configure(text=text)

    def _render_health(self) -> None:
        if "health_water" not in self.panel_widgets:
            return
        health = self.data["health"]
        water: tk.Label = self.panel_widgets["health_water"]  # type: ignore[assignment]
        water_count = max(0, int(health.get("water_count", 0)))
        water.configure(text=f"今天 {water_count} 杯")
        water_visual: tk.Label = self.panel_widgets["health_water_visual"]  # type: ignore[assignment]
        filled = min(water_count, 8)
        water_visual.configure(text=" ".join(["💧"] * filled + ["·"] * (8 - filled)))
        self._render_move_health()
        self._render_medicine_list()

    def _render_move_health(self) -> None:
        if "health_move_timer" not in self.panel_widgets:
            return
        health = self.data["health"]
        interval = max(5, int(health.get("move_interval", 60)))
        try:
            elapsed_seconds = max(0.0, (datetime.now() - datetime.fromisoformat(health.get("last_move", now_iso()))).total_seconds())
        except (TypeError, ValueError):
            elapsed_seconds = 0.0
        elapsed_minutes = int(elapsed_seconds // 60)
        progress = min(100.0, elapsed_seconds / (interval * 60) * 100)
        timer: tk.Label = self.panel_widgets["health_move_timer"]  # type: ignore[assignment]
        if health.get("move_done_date") == date.today().isoformat() and elapsed_seconds < 60:
            timer.configure(text=f"剛完成活動｜下次 {interval} 分鐘後提醒")
        else:
            remaining = max(0, interval - elapsed_minutes)
            timer.configure(text=f"已累計 {elapsed_minutes} 分鐘｜約 {remaining} 分鐘後提醒")
        progress_bar: ttk.Progressbar = self.panel_widgets["health_move_progress"]  # type: ignore[assignment]
        progress_bar.configure(value=progress)

    def _health_animation_tick(self) -> None:
        """每秒更新起身活動累計，提供可見的進度動畫。"""
        try:
            self._render_move_health()
        finally:
            self.root.after(1000, self._health_animation_tick)

    def _set_move_interval(self, value: str) -> None:
        try:
            interval = max(5, int(value))
        except (TypeError, ValueError):
            interval = 60
        self.data["health"]["move_interval"] = interval
        self.health_alerted["move"] = False
        self._save_data()
        self._render_move_health()

    def _medicine_times(self) -> list[str]:
        values = self.data["health"].get("medicine_times", [])
        if not isinstance(values, list):
            return []
        return sorted({str(value) for value in values if re.fullmatch(r"(?:[01]\d|2[0-3]):[0-5]\d", str(value))})

    def _render_medicine_list(self) -> None:
        list_frame = self.panel_widgets.get("medicine_list")
        if not isinstance(list_frame, tk.Frame):
            return
        for child in list_frame.winfo_children():
            child.destroy()
        times = self._medicine_times()
        done = set(self.data["health"].get("medicine_done_times", []))
        if not times:
            tk.Label(list_frame, text="尚未設定服藥時間", bg=self.COLORS["soft"], fg=self.COLORS["muted"], anchor="w", font=("Microsoft JhengHei", 8)).pack(fill="x")
            return
        for med_time in times:
            row = tk.Frame(list_frame, bg=self.COLORS["soft"])
            row.pack(fill="x", pady=1)
            status = "已完成" if med_time in done else "待提醒"
            tk.Label(row, text=f"{med_time}｜{status}", bg=self.COLORS["soft"], fg=self.COLORS["ink"], anchor="w", font=("Microsoft JhengHei", 8, "bold")).pack(side="left", fill="x", expand=True)
            ttk.Button(row, text="已服用", style="Secretary.TButton", command=lambda value=med_time: self._record_health("medicine", value)).pack(side="right")
            ttk.Button(row, text="刪除", style="Secretary.TButton", command=lambda value=med_time: self._remove_medicine_time(value)).pack(side="right", padx=(0, 4))

    def _add_medicine_time(self) -> None:
        value = str(self.panel_widgets["medicine_var"].get()).strip()  # type: ignore[union-attr]
        if not value:
            messagebox.showinfo(APP_NAME, "請先從下拉選單選擇時間。", parent=self.panel)
            return
        times = self._medicine_times()
        if value not in times:
            times.append(value)
            times.sort()
            self.data["health"]["medicine_times"] = times
            self.data["health"]["medicine_time"] = times[0]
            self._save_data()
        self._render_medicine_list()

    def _remove_medicine_time(self, value: str) -> None:
        times = [item for item in self._medicine_times() if item != value]
        health = self.data["health"]
        health["medicine_times"] = times
        health["medicine_time"] = times[0] if times else ""
        health["medicine_done_times"] = [item for item in health.get("medicine_done_times", []) if item in times]
        health["medicine_alerted_times"] = [item for item in health.get("medicine_alerted_times", []) if item in times]
        self._save_data()
        self._render_medicine_list()

    def _record_health(self, kind: str, medicine_time: str | None = None) -> None:
        health = self.data["health"]
        if kind == "water":
            health["water_count"] = int(health.get("water_count", 0)) + 1
            health["last_water"] = now_iso()
            self.health_alerted["water"] = False
            message = "喝水完成！"
        elif kind == "move":
            health["last_move"] = now_iso()
            health["move_done_date"] = date.today().isoformat()
            self.health_alerted["move"] = False
            message = "起身活動完成！"
        else:
            times = self._medicine_times()
            if not times:
                messagebox.showinfo(APP_NAME, "請先從下拉選單新增服藥提醒時間。", parent=self.panel)
                return
            selected = medicine_time if medicine_time in times else times[0]
            done_times = health.setdefault("medicine_done_times", [])
            if selected not in done_times:
                done_times.append(selected)
            health["medicine_done_date"] = date.today().isoformat()
            health["medicine_alerted_times"] = [item for item in health.get("medicine_alerted_times", []) if item != selected]
            self.health_alerted["medicine"] = False
            message = f"{selected} 的服藥紀錄已完成。"
        self._save_data()
        self._render_health()
        self.play("success", 2, "idle", message)

    def _set_medicine_time(self, value: str) -> None:
        value = value.strip()
        if value and not re.fullmatch(r"(?:[01]\d|2[0-3]):[0-5]\d", value):
            messagebox.showwarning(APP_NAME, "請輸入 24 小時格式，例如 13:30。", parent=self.panel)
            return
        if value:
            self.data["health"]["medicine_times"] = sorted(set(self._medicine_times() + [value]))
            self.data["health"]["medicine_time"] = self.data["health"]["medicine_times"][0]
            self.data["health"]["medicine_done_date"] = ""
            self.data["health"]["medicine_done_times"] = []
            self.data["health"]["medicine_alerted_times"] = []
            self.health_alerted["medicine"] = False
            self._save_data()
            self._render_medicine_list()

    def _health_tick(self) -> None:
        health = self.data["health"]
        now = datetime.now()
        if health.get("date") != date.today().isoformat():
            health["date"] = date.today().isoformat()
            health["water_count"] = 0
            health["medicine_done_times"] = []
            health["medicine_alerted_times"] = []
            self._save_data()
            self._render_health()
        if not self.data.get("settings", {}).get("health_reminders", True) or self._in_quiet_hours():
            self.root.after(60_000, self._health_tick)
            return
        alerts: list[str] = []
        move_interval = max(5, int(health.get("move_interval", 60)))
        water_interval = max(5, int(health.get("water_interval", 60)))
        for key, field, minutes, message in [
            ("water", "last_water", water_interval, "忙了一段時間，記得喝口水。"),
            ("move", "last_move", move_interval, f"已經坐了 {move_interval} 分鐘，起身活動一下吧！"),
        ]:
            try:
                elapsed = now - datetime.fromisoformat(health[field])
            except (KeyError, TypeError, ValueError):
                elapsed = timedelta()
            if elapsed >= timedelta(minutes=minutes) and not self.health_alerted[key]:
                self.health_alerted[key] = True
                alerts.append(message)
        current_time = now.strftime("%H:%M")
        done_times = set(health.get("medicine_done_times", []))
        alerted_times = set(health.get("medicine_alerted_times", []))
        for med_time in self._medicine_times():
            if current_time >= med_time and med_time not in done_times and med_time not in alerted_times:
                alerted_times.add(med_time)
                alerts.append(f"服藥提醒時間 {med_time} 到了，請依自己的醫囑處理。")
                break
        health["medicine_alerted_times"] = sorted(alerted_times)
        if alerts:
            self.play("warning", 3, "idle", alerts[0])
        self.root.after(60_000, self._health_tick)

    def _opening_brief(self) -> None:
        if self._in_quiet_hours():
            return
        today = date.today().isoformat()
        active = [task for task in self.data["tasks"] if task.get("status") != "已完成"]
        overdue = [task for task in active if task.get("due_date") and task["due_date"] < today]
        due = [task for task in active if task.get("due_date") == today]
        if overdue:
            self.show_bubble(f"有 {len(overdue)} 件逾期任務，要先看看嗎？", 4200)
        elif due:
            self.show_bubble(f"今天有 {len(due)} 件期限任務。", 3600)

    def _choose_attachments(self) -> None:
        paths = filedialog.askopenfilenames(
            title="交給小綿助整理",
            filetypes=[("支援的檔案", "*.png *.jpg *.jpeg *.webp *.bmp *.tif *.tiff *.pdf *.txt *.md *.csv *.json *.html *.docx *.xlsx *.xlsm *.pptx"), ("所有檔案", "*.*")],
            parent=self.panel,
        )
        self._add_attachment_paths([Path(path) for path in paths])

    def _paste_attachment(self, _event: tk.Event) -> str | None:
        try:
            content = ImageGrab.grabclipboard()
        except OSError:
            return None
        if isinstance(content, Image.Image):
            target = ATTACHMENT_ROOT / f"貼上圖片-{datetime.now().strftime('%Y%m%d-%H%M%S')}.png"
            content.save(target, "PNG")
            self._add_attachment_paths([target])
            return "break"
        if isinstance(content, list):
            paths = [Path(item) for item in content if Path(item).is_file()]
            if paths:
                self._add_attachment_paths(paths)
                return "break"
        return None

    def _add_attachment_paths(self, paths: list[Path]) -> None:
        for path in paths:
            if len(self.attachments) >= 10:
                messagebox.showwarning(APP_NAME, "一次最多加入 10 個附件。", parent=self.panel)
                break
            try:
                if path.stat().st_size > 35 * 1024 * 1024:
                    messagebox.showwarning(APP_NAME, f"{path.name} 超過單檔 35 MB。", parent=self.panel)
                    continue
            except OSError:
                continue
            if path not in self.attachments:
                self.attachments.append(path)
        self._render_attachments()

    def _persist_attachments(self) -> list[Path]:
        """將外部附件複製到桌寵資料夾，避免原檔移動後記事失效。"""
        ATTACHMENT_ROOT.mkdir(parents=True, exist_ok=True)
        stored: list[Path] = []
        for path in self.attachments:
            try:
                if path.resolve().parent == ATTACHMENT_ROOT.resolve():
                    stored.append(path)
                    continue
                safe_name = re.sub(r"[^\w.() -]+", "_", path.name, flags=re.UNICODE)
                target = ATTACHMENT_ROOT / f"{datetime.now().strftime('%Y%m%d-%H%M%S')}-{uuid.uuid4().hex[:6]}-{safe_name}"
                shutil.copy2(path, target)
                stored.append(target)
            except OSError:
                stored.append(path)
        self.attachments = stored
        return stored

    def _extract_attachment_text(self, path: Path) -> str:
        cache_key = str(path.resolve())
        if cache_key in self.attachment_content_cache:
            return self.attachment_content_cache[cache_key]
        suffix = path.suffix.lower()
        text = ""
        try:
            if suffix in {".txt", ".md", ".csv", ".json", ".html", ".htm"}:
                text = path.read_text(encoding="utf-8", errors="replace")
                if suffix in {".html", ".htm"}:
                    text = re.sub(r"<[^>]+>", " ", text)
            elif suffix == ".pdf":
                from pypdf import PdfReader

                reader = PdfReader(str(path))
                text = "\n".join((page.extract_text() or "") for page in reader.pages[:30])
            elif suffix == ".docx":
                from docx import Document

                document = Document(str(path))
                parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
                for table in document.tables:
                    parts.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
                text = "\n".join(parts)
            elif suffix in {".xlsx", ".xlsm"}:
                from openpyxl import load_workbook

                workbook = load_workbook(path, read_only=True, data_only=True)
                parts: list[str] = []
                for sheet in workbook.worksheets[:10]:
                    parts.append(f"工作表：{sheet.title}")
                    for row in sheet.iter_rows(max_row=300, max_col=30, values_only=True):
                        values = [str(value) for value in row if value not in (None, "")]
                        if values:
                            parts.append(" | ".join(values))
                workbook.close()
                text = "\n".join(parts)
            elif suffix == ".pptx":
                from pptx import Presentation

                presentation = Presentation(str(path))
                parts = []
                for index, slide in enumerate(presentation.slides[:80], start=1):
                    slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                    if slide_text:
                        parts.append(f"第 {index} 張：" + "\n".join(slide_text))
                text = "\n".join(parts)
            elif suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}:
                text = self._windows_ocr(path)
            else:
                text = ""
        except (OSError, ValueError, ImportError, RuntimeError) as exc:
            text = f"[附件讀取失敗：{exc}]"
        text = re.sub(r"[ \t]+", " ", text).strip()[:20_000]
        self.attachment_content_cache[cache_key] = text
        return text

    def _windows_ocr(self, path: Path) -> str:
        if sys.platform != "win32":
            return "[圖片文字辨識僅支援 Windows 10／11]"
        script = Path(__file__).resolve().parent / "windows_ocr.ps1"
        if not script.exists():
            return "[缺少 Windows 圖片辨識元件]"
        creation_flags = getattr(subprocess, "CREATE_NO_WINDOW", 0)
        try:
            result = subprocess.run(
                ["powershell.exe", "-NoProfile", "-NonInteractive", "-File", str(script), str(path.resolve())],
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=45,
                creationflags=creation_flags,
                check=False,
            )
        except (OSError, subprocess.TimeoutExpired):
            return "[Windows 圖片文字辨識沒有回應]"
        recognized = result.stdout.strip()
        if recognized:
            return recognized
        return "[圖片中沒有辨識到文字]"

    def _collect_attachment_content(self) -> tuple[str, list[str]]:
        parts: list[str] = []
        statuses: list[str] = []
        for path in self.attachments:
            content = self._extract_attachment_text(path)
            readable = bool(content and not content.startswith("["))
            statuses.append(f"{path.name}：{'已讀取內容' if readable else content or '不支援此格式'}")
            if content:
                parts.append(f"【附件：{path.name}】\n{content}")
        return "\n\n".join(parts)[:40_000], statuses

    def _render_attachments(self) -> None:
        if "attachments" not in self.panel_widgets:
            return
        label: tk.Label = self.panel_widgets["attachments"]  # type: ignore[assignment]
        label.configure(text="尚未加入附件" if not self.attachments else "待辨識附件：" + "、".join(path.name for path in self.attachments))

    def _clear_attachments(self) -> None:
        self.attachments = []
        self.attachment_content_cache = {}
        self._render_attachments()

    def _voice_handoff(self) -> None:
        handoff: tk.Text = self.panel_widgets["handoff"]  # type: ignore[assignment]
        handoff.focus_set()
        self.panel.update_idletasks()
        if sys.platform == "win32":
            try:
                import ctypes
                user32 = ctypes.windll.user32
                keyup = 0x0002
                user32.keybd_event(0x5B, 0, 0, 0)
                user32.keybd_event(ord("H"), 0, 0, 0)
                user32.keybd_event(ord("H"), 0, keyup, 0)
                user32.keybd_event(0x5B, 0, keyup, 0)
            except (AttributeError, OSError):
                pass
        self.show_bubble("已開啟 Windows 語音輸入；說完後再交給我整理。", 3000)

    def _analyze_handoff(self) -> None:
        handoff: tk.Text = self.panel_widgets["handoff"]  # type: ignore[assignment]
        text = handoff.get("1.0", "end").strip()
        if not text and not self.attachments:
            messagebox.showinfo(APP_NAME, "請輸入交代內容，或加入圖片／檔案。", parent=self.panel)
            return
        if self.attachments:
            self.show_bubble("正在讀取附件內容……", 3000)
            if self.panel:
                self.panel.update_idletasks()
            self._persist_attachments()
        attachment_text, attachment_statuses = self._collect_attachment_content()
        source = "\n\n".join(part for part in (text, attachment_text) if part).strip()
        if not source:
            source = "處理附件：" + "、".join(path.name for path in self.attachments)
        title_source = text or attachment_text[:800] or source
        due_date = self._parse_date(source)
        due_time_match = re.search(r"(?:上午|下午|晚上)?\s*(\d{1,2})[:：點](\d{2})", source)
        due_time = f"{int(due_time_match.group(1)):02d}:{due_time_match.group(2)}" if due_time_match else ""
        title = re.sub(r"(今天|明天|後天|請|幫我|記得|提醒我|提醒|加入|新增|待辦)", " ", title_source)
        title = re.sub(r"\s+", " ", title).strip(" ，。,.！!")[:80] or "未命名事項"
        kind = "行程" if re.search(r"會議|開會|日曆|行事曆", source) else "任務"
        waiting_for = "" if not re.search(r"等待|回覆|收齊|調查", source) else "待確認對象"
        self.draft = {
            "id": str(uuid.uuid4()), "title": title, "kind": kind, "due_date": due_date,
            "due_time": due_time, "priority": "高" if due_date <= (date.today() + timedelta(days=1)).isoformat() else "中",
            "status": "等待回覆" if waiting_for else "未開始", "waiting_for": waiting_for,
            "attachment_names": [path.name for path in self.attachments],
            "attachment_paths": [str(path) for path in self.attachments],
            "attachment_statuses": attachment_statuses,
            "source": source,
            "created_at": now_iso(),
        }
        label: tk.Label = self.panel_widgets["draft_label"]  # type: ignore[assignment]
        readable_count = sum("已讀取內容" in status for status in attachment_statuses)
        label.configure(text=f"建議建立{kind}：{title}\n日期：{due_date} {due_time or '未指定時間'}｜優先級：{self.draft['priority']}\n附件：{len(self.attachments)} 個，已讀取 {readable_count} 個內容。請確認後才會寫入本機任務清單。")
        frame: tk.Frame = self.panel_widgets["draft_frame"]  # type: ignore[assignment]
        frame.pack(fill="x", pady=(10, 0))
        frame.lift()
        self.play("think", 2, "idle", "我整理好了，請你確認。")

    @staticmethod
    def _parse_date(text: str) -> str:
        explicit = re.search(r"(20\d{2})[年/-](\d{1,2})[月/-](\d{1,2})", text)
        if explicit:
            return date(int(explicit.group(1)), int(explicit.group(2)), int(explicit.group(3))).isoformat()
        target = date.today()
        if "後天" in text:
            target += timedelta(days=2)
        elif "明天" in text:
            target += timedelta(days=1)
        return target.isoformat()

    def _confirm_draft(self) -> None:
        if not self.draft:
            return
        self.data["tasks"].append(self.draft)
        self._save_data()
        self.draft = None
        handoff: tk.Text = self.panel_widgets["handoff"]  # type: ignore[assignment]
        handoff.delete("1.0", "end")
        self.attachments = []
        self.attachment_content_cache = {}
        self._render_attachments()
        frame: tk.Frame = self.panel_widgets["draft_frame"]  # type: ignore[assignment]
        frame.pack_forget()
        self._render_dashboard()
        self.play("success", 2, "idle", "完成！已加入桌面秘書任務清單。")

    def _save_note(self) -> None:
        handoff: tk.Text = self.panel_widgets["handoff"]  # type: ignore[assignment]
        text = handoff.get("1.0", "end").strip()
        if not text and not self.attachments:
            messagebox.showinfo(APP_NAME, "目前沒有可以儲存的內容。", parent=self.panel)
            return
        if self.attachments:
            self._persist_attachments()
        attachment_text, attachment_statuses = self._collect_attachment_content()
        self.data["notes"].insert(0, {
            "text": text or attachment_text[:500] or "附件記事",
            "attachments": [path.name for path in self.attachments],
            "attachment_paths": [str(path) for path in self.attachments],
            "attachment_text": attachment_text,
            "attachment_statuses": attachment_statuses,
            "created_at": now_iso(),
        })
        self._save_data()
        handoff.delete("1.0", "end")
        self.attachments = []
        self.attachment_content_cache = {}
        self._render_attachments()
        self._render_notes()
        self.play("success", 2, "idle", "記事已保存在本機。")


if __name__ == "__main__":
    app = SecretaryPet()
    if "--open" in sys.argv:
        app.root.after(500, app.open_secretary)
    app.run()
